from pptx import Presentation
from pptx.util import Inches
import numpy as np
from scipy.io.wavfile import write

def create_test_presentation():
    """Creates a simple 2-slide presentation with speaker notes."""
    prs = Presentation()

    # Slide 1
    slide1_layout = prs.slide_layouts[5]  # Title only layout
    slide1 = prs.slides.add_slide(slide1_layout)
    title1 = slide1.shapes.title
    title1.text = "Slide 1: The Journey Begins"
    notes_slide1 = slide1.notes_slide
    notes_slide1.notes_text_frame.text = "This is the first slide. Welcome to our presentation on AI-driven video generation."

    # Slide 2
    slide2_layout = prs.slide_layouts[5]
    slide2 = prs.slides.add_slide(slide2_layout)
    title2 = slide2.shapes.title
    title2.text = "Slide 2: The Grand Finale"
    notes_slide2 = slide2.notes_slide
    notes_slide2.notes_text_frame.text = "This is the second and final slide. Thank you for your attention."

    prs.save("test.pptx")
    print("Created test.pptx")

def create_test_wav():
    """Creates a short, silent WAV file for voice cloning tests."""
    samplerate = 24000
    duration_s = 1
    # Create a silent audio signal
    amplitude = np.iinfo(np.int16).max * 0.0 # Zero amplitude
    data = amplitude * np.sin(2. * np.pi * 440. * np.arange(samplerate * duration_s) / samplerate)
    write("test.wav", samplerate, data.astype(np.int16))
    print("Created test.wav")

if __name__ == "__main__":
    create_test_presentation()
    create_test_wav()
