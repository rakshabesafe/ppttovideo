from app.workers.celery_app_cpu import app as celery_app
from app.db.session import SessionLocal
from app import crud
from app.services.minio_service import minio_service
from pptx import Presentation
import requests
import io
import os
import tempfile
from celery import chord, group
from moviepy import ImageClip, AudioFileClip, concatenate_videoclips
# Import tasks through celery to avoid direct module imports
# This prevents MoviePy being loaded in GPU worker

@celery_app.task(name="app.workers.tasks_cpu.decompose_presentation")
def decompose_presentation(job_id: int):
    # ... (previous implementation is correct)
    db = SessionLocal()
    try:
        job = crud.get_presentation_job(db, job_id)
        if not job:
            print(f"Job {job_id} not found.")
            return

        # Create decomposition task tracking
        decomp_task = crud.create_job_task(db, job_id, "decomposition")
        crud.update_task_status(db, task_id=decomp_task.id, status="running", progress_message="Extracting slides from presentation")
        
        crud.update_job_status(db, job_id, "processing_slides", current_stage="decomposing")

        pptx_object_name = job.s3_pptx_path.split('/', 2)[-1]
        pptx_bucket_name = job.s3_pptx_path.split('/')[1]

        response = minio_service.client.get_object(pptx_bucket_name, pptx_object_name)
        pptx_data = io.BytesIO(response.read())
        response.close()
        response.release_conn()

        prs = Presentation(pptx_data)
        num_slides = len(prs.slides)
        
        # Update job with slide count
        crud.update_job_slides(db, job_id, num_slides)
        crud.update_task_status(db, task_id=decomp_task.id, progress_message=f"Processing {num_slides} slides")
        
        notes_paths = []
        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            note_object_name = f"{job_id}/notes/slide_{slide_number}.txt"
            minio_service.upload_file(
                bucket_name="presentations",
                object_name=note_object_name,
                data=io.BytesIO(notes.encode('utf-8')),
                length=len(notes.encode('utf-8'))
            )
            notes_paths.append(f"/presentations/{note_object_name}")

        libreoffice_url = "http://libreoffice:8100/convert"
        payload = {"bucket_name": pptx_bucket_name, "object_name": pptx_object_name}
        res = requests.post(libreoffice_url, json=payload)
        res.raise_for_status()
        image_paths = res.json().get("image_paths", [])

        if len(image_paths) != num_slides:
            raise Exception(f"Mismatch between number of images ({len(image_paths)}) and slides ({num_slides}).")

        # Complete decomposition task
        crud.update_task_status(db, task_id=decomp_task.id, status="completed", progress_message=f"Successfully processed {num_slides} slides")
        
        # Create audio synthesis tasks for all slides and collect their AsyncResult objects
        audio_task_results = []
        for i in range(num_slides):
            # Create task tracking record
            audio_task_db = crud.create_job_task(db, job_id, "audio_synthesis", slide_number=i+1)
            
            # Send audio synthesis task to GPU worker
            task_id = f"synthesize_audio_{job_id}_{i+1}"
            print(f"Sending audio synthesis task for job {job_id}, slide {i+1}")
            result = celery_app.send_task('app.workers.tasks_gpu.synthesize_audio', 
                                        args=(job_id, i+1), 
                                        task_id=task_id,
                                        queue='gpu_tasks')
            print(f"Task sent with result: {result}")
            
            # Update task with Celery task ID
            crud.update_task_status(db, task_id=audio_task_db.id, set_celery_task_id=str(result.id))
            
            audio_task_results.append(result)
        
        print(f"Image paths for job {job_id}: {image_paths}")
        print(f"Created {len(audio_task_results)} audio synthesis tasks")
        
        # Schedule video assembly with dependency tracking to wait for audio synthesis completion
        print(f"Scheduling video assembly for job {job_id} with dependency tracking")
        assemble_video_with_deps.apply_async(args=(image_paths, job_id, [result.id for result in audio_task_results]), queue='cpu_tasks')

        crud.update_job_status(db, job_id, "synthesizing_audio", current_stage="synthesizing_audio")

    except Exception as e:
        crud.update_job_status(db, job_id, "failed")
        print(f"Error in decompose_presentation for job {job_id}: {e}")
    finally:
        db.close()


@celery_app.task(name="app.workers.tasks_cpu.assemble_video_with_deps")
def assemble_video_with_deps(image_paths_from_libreoffice, job_id: int, audio_task_ids: list, max_wait_time: int = 600):
    """
    Assembles video after waiting for all audio synthesis tasks to complete.
    
    Args:
        image_paths_from_libreoffice: List of image paths returned by LibreOffice service
        job_id: The presentation job ID
        audio_task_ids: List of audio synthesis task IDs to wait for
        max_wait_time: Maximum time to wait for audio tasks (default: 10 minutes)
    """
    import time
    from celery import current_app
    
    db = SessionLocal()
    try:
        print(f"Waiting for audio synthesis tasks to complete for job {job_id}")
        print(f"Audio task IDs: {audio_task_ids}")
        print(f"Maximum wait time: {max_wait_time} seconds")
        
        start_time = time.time()
        completed_tasks = set()
        
        # Wait for all audio tasks to complete
        while len(completed_tasks) < len(audio_task_ids):
            elapsed_time = time.time() - start_time
            
            # Check timeout
            if elapsed_time > max_wait_time:
                print(f"Timeout reached ({max_wait_time}s). Proceeding with available audio files.")
                crud.update_job_status(db, job_id, "failed", error_message=f"Audio synthesis timeout after {max_wait_time}s")
                return
            
            # Check each task status
            for task_id in audio_task_ids:
                if task_id not in completed_tasks:
                    result = current_app.AsyncResult(task_id)
                    if result.ready():
                        if result.successful():
                            print(f"Audio task {task_id} completed successfully")
                            completed_tasks.add(task_id)
                        else:
                            print(f"Audio task {task_id} failed: {result.result}")
                            # Continue anyway - we'll handle missing audio in assemble_video
                            completed_tasks.add(task_id)
            
            # If not all complete, wait a bit and check again
            if len(completed_tasks) < len(audio_task_ids):
                remaining = len(audio_task_ids) - len(completed_tasks)
                print(f"Waiting for {remaining} audio tasks to complete... ({elapsed_time:.1f}s elapsed)")
                time.sleep(10)  # Check every 10 seconds
        
        total_time = time.time() - start_time
        print(f"All audio synthesis tasks completed in {total_time:.1f} seconds")
        
        # Now call the original video assembly function
        return assemble_video(image_paths_from_libreoffice, job_id)
        
    except Exception as e:
        crud.update_job_status(db, job_id, "failed")
        print(f"Error in assemble_video_with_deps for job {job_id}: {e}")
        raise
    finally:
        db.close()


@celery_app.task(name="app.workers.tasks_cpu.assemble_video")
def assemble_video(image_paths_from_libreoffice, job_id: int):
    """
    Assembles the final video from slide images and synthesized audio.
    Args:
        image_paths_from_libreoffice: List of image paths returned by LibreOffice service
        job_id: The presentation job ID
    """
    db = SessionLocal()
    try:
        crud.update_job_status(db, job_id, "assembling_video")

        with tempfile.TemporaryDirectory() as temp_dir:
            # 1. Download images using paths from LibreOffice service
            image_local_paths = {}
            print(f"Processing {len(image_paths_from_libreoffice)} images from LibreOffice")
            
            for i, image_s3_path in enumerate(image_paths_from_libreoffice):
                slide_num = i + 1
                # Remove leading slash if present: /presentations/uuid/images/slide-1.png -> presentations/uuid/images/slide-1.png
                clean_path = image_s3_path.lstrip('/')
                # Extract bucket and object: presentations/uuid/images/slide-1.png -> bucket='presentations', object='uuid/images/slide-1.png'  
                if clean_path.startswith('presentations/'):
                    object_name = clean_path[len('presentations/'):]
                    bucket_name = 'presentations'
                else:
                    # Fallback - assume it's the full object name
                    bucket_name = 'presentations' 
                    object_name = clean_path
                    
                local_path = os.path.join(temp_dir, f"slide_{slide_num}.png")
                print(f"Downloading image {slide_num}: {bucket_name}/{object_name} -> {local_path}")
                minio_service.client.fget_object(bucket_name, object_name, local_path)
                image_local_paths[slide_num] = local_path

            # 2. Download audio files
            audio_files = minio_service.client.list_objects("presentations", prefix=f"{job_id}/audio/")
            audio_paths = {}

            for aud in audio_files:
                local_path = os.path.join(temp_dir, os.path.basename(aud.object_name))
                minio_service.client.fget_object("presentations", aud.object_name, local_path)
                # slide_1.wav -> 1
                slide_num = int(os.path.splitext(os.path.basename(local_path))[0].split('_')[1])
                audio_paths[slide_num] = local_path
                print(f"Downloaded audio {slide_num}: {local_path}")

            # 3. Create video clips
            print(f"Creating video clips for {len(image_local_paths)} slides")
            clips = []
            for i in sorted(image_local_paths.keys()):
                image_path = image_local_paths[i]
                audio_path = audio_paths.get(i)

                if not audio_path:
                    raise Exception(f"Missing audio for slide {i}")

                try:
                    audio_clip = AudioFileClip(audio_path)
                    # Try different approaches for creating image clip
                    try:
                        # Method 1: Constructor with duration
                        image_clip = ImageClip(image_path, duration=audio_clip.duration)
                    except:
                        try:
                            # Method 2: Create then set duration
                            image_clip = ImageClip(image_path).set_duration(audio_clip.duration)
                        except:
                            # Method 3: Use duration parameter differently
                            image_clip = ImageClip(image_path)
                            image_clip.duration = audio_clip.duration
                    
                    # Set fps if possible
                    try:
                        image_clip.fps = 24
                    except:
                        pass  # Ignore fps setting if not supported
                    
                    final_clip = image_clip.with_audio(audio_clip)
                    clips.append(final_clip)
                    print(f"Successfully created clip for slide {i}")
                except Exception as e:
                    print(f"Error creating clip for slide {i}: {e}")
                    # Create a minimal clip as fallback
                    audio_clip = AudioFileClip(audio_path)
                    # Use a simple approach - create video clip without image
                    from moviepy import ColorClip
                    color_clip = ColorClip(size=(1920,1080), color=(0,0,0), duration=audio_clip.duration)
                    final_clip = color_clip.with_audio(audio_clip)
                    clips.append(final_clip)
                    print(f"Created fallback black screen for slide {i}")

            # 3. Concatenate and write final video
            final_video = concatenate_videoclips(clips)
            output_filename = f"{job_id}.mp4"
            local_output_path = os.path.join(temp_dir, output_filename)
            final_video.write_videofile(local_output_path, codec='libx264', audio_codec='aac')

            # 4. Upload to MinIO output bucket
            with open(local_output_path, "rb") as f:
                s3_path = minio_service.upload_file(
                    bucket_name="output",
                    object_name=output_filename,
                    data=f,
                    length=os.path.getsize(local_output_path)
                )

            # 5. Update job status in DB
            crud.update_job_status(db, job_id, "completed", video_path=s3_path)

    except Exception as e:
        crud.update_job_status(db, job_id, "failed")
        print(f"Error in assemble_video for job {job_id}: {e}")
    finally:
        db.close()
